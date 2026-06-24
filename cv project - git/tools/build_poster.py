"""
build_poster.py — fills the poster template with this project's content.

1. Generates a vertical pipeline diagram (outputs/pipeline.png) for Methodology.
2. Opens 'Poster Template-new.pptx', replaces the sample text / tables / diagram
   with Fire & Smoke Vision System content, and saves 'Fire_Smoke_Poster.pptx'.

Run:  python tools/build_poster.py
Personal fields (name / reg no / course) are left as [PLACEHOLDERS] — edit them
in PowerPoint, or change the constants below and re-run.
"""
import os

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor

HERE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATE = os.path.join(HERE, "Poster Template-new.pptx")
OUT_PPTX = os.path.join(HERE, "Fire_Smoke_Poster.pptx")
PIPE_PNG = os.path.join(HERE, "outputs", "pipeline.png")

CLAY = "#CC785C"
INK = "#1A1916"

# ---- personal / header fields (EDIT THESE) --------------------------------
MEMBERS = "Salman Naseem (2023-CS-78)        Huzaifa Arshad (2023-CS-86)"
COURSE = "Computer Vision"
INSTRUCTOR = "Muhammad Waseem"
SEG_IMG = os.path.join(HERE, "outputs", "seg_fire.jpg")
CM_IMG = os.path.join(HERE, "outputs", "confusion_matrix.png")
BODY = 35   # body font size requested

TITLE = ("Fire & Smoke Vision System: A Multi-Module Deep Learning Approach to "
         "Fire Detection, Localization & Severity Assessment")


# =========================================================================
# 1. Pipeline diagram
# =========================================================================
def make_pipeline():
    boxes = [
        ("Dataset", "5,000 images · LabelMe polygons\n→ classification + detection + segmentation labels", "#EFE7DC"),
        ("Input Image", "JPG / PNG", "#FFFFFF"),
        ("Preprocessing", "Resize · Normalize · Augment", "#FFFFFF"),
        ("Module 1 — Classification", "ResNet18  →  fire / normal", CLAY),
        ("Module 2 — Detection", "YOLOv8  →  bounding boxes", CLAY),
        ("Module 3 — Segmentation", "YOLOv8-seg  →  pixel masks", CLAY),
        ("Module 4 — Explainability", "Grad-CAM heatmap", CLAY),
        ("Module 5 — Severity", "% frame on fire → risk level", CLAY),
        ("React Web Application", "FastAPI backend serves all modules", "#EFE7DC"),
    ]
    n = len(boxes)
    fig, ax = plt.subplots(figsize=(6.3, 12.7))
    ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.axis("off")
    top, bot = 0.97, 0.03
    gap = (top - bot) / n
    bh = gap * 0.62
    centers = []
    for i, (head, sub, fill) in enumerate(boxes):
        y = top - gap * (i + 0.5)
        centers.append(y)
        dark = fill == CLAY
        ax.add_patch(FancyBboxPatch((0.06, y - bh / 2), 0.88, bh,
                                    boxstyle="round,pad=0.012,rounding_size=0.02",
                                    linewidth=2, edgecolor=CLAY,
                                    facecolor=fill, mutation_aspect=0.4))
        ax.text(0.5, y + bh * 0.16, head, ha="center", va="center",
                fontsize=13, fontweight="bold",
                color="white" if dark else INK)
        ax.text(0.5, y - bh * 0.22, sub, ha="center", va="center",
                fontsize=10, color="#F3E7E1" if dark else "#5A574F")
    for i in range(n - 1):
        ax.add_patch(FancyArrowPatch((0.5, centers[i] - bh / 2),
                                     (0.5, centers[i + 1] + bh / 2),
                                     arrowstyle="-|>", mutation_scale=16,
                                     linewidth=1.6, color="#9A958C"))
    fig.savefig(PIPE_PNG, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    print("wrote", PIPE_PNG)


# =========================================================================
# 2. Fill the template
# =========================================================================
def set_text(shape, lines, size, bold=False, color=INK):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold
        r.font.name = "Calibri"
        r.font.color.rgb = RGBColor.from_string(color.lstrip("#"))
    return shape


def fill_table(table, rows, header_size=22, body_size=20):
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            if ri >= len(table.rows) or ci >= len(table.columns):
                continue
            cell = table.cell(ri, ci)
            cell.text = val
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(header_size if ri == 0 else body_size)
                    run.font.bold = ri == 0
                    run.font.name = "Calibri"


def starts(shape, prefix):
    return shape.has_text_frame and shape.text_frame.text.strip().startswith(prefix)


def main():
    make_pipeline()
    prs = Presentation(TEMPLATE)
    slide = prs.slides[0]

    abstract = [
        "A multi-module computer-vision system that detects fire and smoke in images. "
        "It unifies five deep-learning tasks — classification (ResNet18), detection and "
        "segmentation (YOLOv8), Grad-CAM explainability and a severity score — trained "
        "from one dataset of 5,000 images. On unseen test data it reaches 99.0% "
        "accuracy, 0.79 detection mAP and 0.79 mask mAP.",
    ]
    introduction = [
        "Fire and smoke threaten life and property; early visual detection is vital for "
        "surveillance and warning systems. Most solutions only classify 'fire or not' "
        "and cannot show where the fire is or how severe it is.",
        "",
        "Objectives:",
        "•  Add localization, segmentation and risk scoring.",
        "•  Train all tasks from one dataset.",
        "•  Explain predictions with Grad-CAM.",
        "•  Provide an easy web interface.",
    ]
    results_text = [
        "Evaluated on a held-out test set of 500 unseen images. Test scores match or "
        "beat validation, so the models generalize without overfitting. Smoke localizes "
        "slightly better than fire.",
    ]
    conclusion = [
        "One dataset powers five fire-analysis tasks in a single explainable web tool, "
        "reaching 99% accuracy and ~0.79 mAP. Future work: higher recall, real-time "
        "video, and edge deployment.",
    ]
    references = [
        "[1] Redmon et al. YOLO. CVPR 2016.",
        "[2] Jocher et al. Ultralytics YOLOv8. 2023.",
        "[3] He et al. Deep Residual Learning (ResNet). CVPR 2016.",
        "[4] Selvaraju et al. Grad-CAM. ICCV 2017.",
        "[5] Muhammad et al. Fire Detection with CNNs. IEEE Access 2018.",
    ]
    related = [
        ["Study", "Focus", "Strength", "Limitation"],
        ["Muhammad+ 2018", "CNN fire classification", "High accuracy", "No localization"],
        ["Li & Zhao 2020", "Faster R-CNN fire", "Bounding boxes", "Heavy; no masks"],
        ["Jadon+ 2019 (FireNet)", "Lightweight CNN", "Real-time", "Classify only"],
        ["Shamsoshoara 2021 (FLAME)", "Aerial fire seg.", "Pixel masks", "Single task"],
        ["YOLOv8 (Ultralytics) 2023", "General det./seg.", "Fast, SOTA", "Not fire-specific"],
        ["This work, 2025", "5 unified modules + XAI", "Explainable, end-to-end", "Single-domain data"],
    ]
    results_table = [
        ["Module", "Metric", "Score"],
        ["Classification (ResNet18)", "Test accuracy", "99.0%"],
        ["Detection (YOLOv8)", "mAP@50", "0.787"],
        ["Segmentation (YOLOv8-seg)", "Mask mAP@50", "0.792"],
        ["Detection per-class", "fire / smoke mAP@50", "0.75 / 0.80"],
        ["Severity", "Frame coverage", "NONE→CRITICAL"],
    ]

    # --- text replacements (match by current sample text) ---
    for sh in slide.shapes:
        if starts(sh, "Deep Research Horizon"):
            set_text(sh, [TITLE], 46, bold=True, color="#FFFFFF")
        elif starts(sh, "Student Name:"):
            set_text(sh, [MEMBERS,
                          f"Course: {COURSE}      Instructor: {INSTRUCTOR}"], 30, color="#FFFFFF")
        elif starts(sh, "This study introduces"):
            set_text(sh, abstract, BODY)
        elif starts(sh, "The complexity of research"):
            set_text(sh, introduction, BODY)
        elif starts(sh, "Evaluated on topics"):
            set_text(sh, results_text, BODY)
        elif starts(sh, "The system automates IEEE"):
            set_text(sh, conclusion, BODY)
        elif starts(sh, "Liu et al., 2024"):
            set_text(sh, references, 30)

    # --- tables ---
    results_tbl_shape = None
    for sh in slide.shapes:
        if sh.has_table:
            t = sh.table
            if len(t.columns) == 4:
                fill_table(t, related)
            elif len(t.columns) == 3:
                fill_table(t, results_table)
                results_tbl_shape = sh

    # --- insert two seeable result images side by side, above the table ----
    from PIL import Image as _Img

    def place_by_height(path, x_in, y_in, h_in):
        iw, ih = _Img.open(path).size
        w_in = h_in * iw / ih
        slide.shapes.add_picture(path, Inches(x_in), Inches(y_in),
                                 Inches(w_in), Inches(h_in))
        return w_in

    if results_tbl_shape is not None:
        y_img = 23.2
        h_img = 5.0
        w1 = place_by_height(SEG_IMG, 18.0, y_img, h_img)        # segmentation result
        x2 = 18.0 + w1 + 0.5
        w2 = place_by_height(CM_IMG, x2, y_img, h_img)           # confusion matrix
        cap = slide.shapes.add_textbox(Inches(18.0), Inches(y_img + h_img + 0.05),
                                       Inches(16.4), Inches(0.6))
        set_text(cap, ["Segmentation + severity result (left) · Classifier confusion matrix (right)"],
                 20, color="#5A574F")
        # move metrics table below the images and tighten rows
        results_tbl_shape.top = Inches(y_img + h_img + 0.9)
        for row in results_tbl_shape.table.rows:
            row.height = Inches(0.8)

    # --- methodology diagram: replace 'Picture 2' ---
    for sh in list(slide.shapes):
        if sh.name == "Picture 2":
            l, t, w, h = sh.left, sh.top, sh.width, sh.height
            sh._element.getparent().remove(sh._element)
            slide.shapes.add_picture(PIPE_PNG, l, t, w, h)
            print("replaced methodology diagram")
            break

    try:
        prs.save(OUT_PPTX)
        print("\nSaved poster ->", OUT_PPTX)
    except PermissionError:
        alt = OUT_PPTX.replace(".pptx", "_NEW.pptx")
        prs.save(alt)
        print("\n[!] Fire_Smoke_Poster.pptx is open in PowerPoint — saved to", alt,
              "\n    Close PowerPoint and rename it, or re-run this script.")


if __name__ == "__main__":
    main()
